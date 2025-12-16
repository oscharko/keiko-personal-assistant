#!/usr/bin/env python3
"""
Convert PowerPoint presentation to PDF including speaker notes.

This script converts a PPTX file to a PDF document where each page contains
the slide image (rendered via LibreOffice) followed by the speaker notes below it.

Requirements:
    - python-pptx: For extracting speaker notes from PPTX files
    - reportlab: For PDF creation
    - Pillow: For image processing
    - PyMuPDF: For PDF to image conversion
    - LibreOffice: For rendering slides to images (must be installed separately)

Installation:
    pip install python-pptx reportlab Pillow PyMuPDF
    brew install --cask libreoffice  # macOS
"""

import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Optional

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from reportlab.lib.colors import HexColor
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import (
    Image as RLImage,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
)


def find_libreoffice() -> Optional[Path]:
    """
    Find the LibreOffice executable on the system.

    Returns:
        Path to the LibreOffice executable or None if not found.
    """
    # Common LibreOffice paths on different systems
    possible_paths = [
        # macOS
        Path("/Applications/LibreOffice.app/Contents/MacOS/soffice"),
        # Linux
        Path("/usr/bin/soffice"),
        Path("/usr/bin/libreoffice"),
        Path("/usr/local/bin/soffice"),
        # Windows
        Path("C:/Program Files/LibreOffice/program/soffice.exe"),
        Path("C:/Program Files (x86)/LibreOffice/program/soffice.exe"),
    ]

    for path in possible_paths:
        if path.exists():
            return path

    # Try to find via PATH
    soffice = shutil.which("soffice")
    if soffice:
        return Path(soffice)

    libreoffice = shutil.which("libreoffice")
    if libreoffice:
        return Path(libreoffice)

    return None


def export_slides_as_images(
        pptx_path: Path,
        output_dir: Path,
        libreoffice_path: Path,
) -> list[Path]:
    """
    Export PowerPoint slides as PNG images using LibreOffice.

    Args:
        pptx_path: Path to the PPTX file.
        output_dir: Directory where images will be saved.
        libreoffice_path: Path to the LibreOffice executable.

    Returns:
        List of paths to the generated PNG images, sorted by slide number.
    """
    # First convert PPTX to PDF using LibreOffice
    cmd = [
        str(libreoffice_path),
        "--headless",
        "--convert-to", "pdf",
        "--outdir", str(output_dir),
        str(pptx_path),
    ]

    result = subprocess.run(
        cmd,
        capture_output=True,
        text=True,
        timeout=300,  # 5 minutes timeout
    )

    if result.returncode != 0:
        raise RuntimeError(
            f"LibreOffice conversion failed: {result.stderr}"
        )

    # Find the generated PDF
    pdf_path = output_dir / f"{pptx_path.stem}.pdf"
    if not pdf_path.exists():
        raise FileNotFoundError(
            f"Expected PDF not found: {pdf_path}"
        )

    # Convert PDF pages to images using sips (macOS) or other tools
    # For cross-platform compatibility, we'll use a different approach
    # We'll keep the PDF and extract pages using reportlab
    return [pdf_path]


def extract_speaker_notes(pptx_path: Path) -> list[str]:
    """
    Extract speaker notes from all slides in a PowerPoint presentation.

    Args:
        pptx_path: Path to the PPTX file.

    Returns:
        List of speaker notes, one per slide.
    """
    prs = Presentation(str(pptx_path))
    notes = []

    for slide in prs.slides:
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            notes.append(notes_text)
        else:
            notes.append("")

    return notes


def get_slide_count(pptx_path: Path) -> int:
    """
    Get the number of slides in a PowerPoint presentation.

    Args:
        pptx_path: Path to the PPTX file.

    Returns:
        Number of slides.
    """
    prs = Presentation(str(pptx_path))
    return len(prs.slides)


def escape_xml(text: str) -> str:
    """
    Escape special XML characters in text.

    Args:
        text: The text to escape.

    Returns:
        Escaped text safe for XML/HTML.
    """
    return (
        text
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
    )


def create_pdf_with_notes(
        slide_images: list[Path],
        notes: list[str],
        output_path: Path,
) -> None:
    """
    Create a PDF with slide images and speaker notes.

    Each page contains:
    1. Slide image (visual representation)
    2. Speaker notes below the image

    Args:
        slide_images: List of paths to slide images (PNG files).
        notes: List of speaker notes corresponding to each slide.
        output_path: Path where the output PDF will be saved.
    """
    # Create the PDF document
    doc = SimpleDocTemplate(
        str(output_path),
        pagesize=A4,
        rightMargin=1 * cm,
        leftMargin=1 * cm,
        topMargin=1 * cm,
        bottomMargin=1 * cm,
    )

    # Calculate available dimensions
    page_width = A4[0] - 2 * cm

    # Define styles
    styles = getSampleStyleSheet()

    slide_header_style = ParagraphStyle(
        "SlideHeader",
        parent=styles["Heading1"],
        fontSize=14,
        textColor=HexColor("#2E4057"),
        spaceBefore=0,
        spaceAfter=8,
        fontName="Helvetica-Bold",
    )

    notes_header_style = ParagraphStyle(
        "NotesHeader",
        parent=styles["Normal"],
        fontSize=10,
        spaceBefore=10,
        spaceAfter=4,
        textColor=HexColor("#2E4057"),
        fontName="Helvetica-Bold",
    )

    notes_style = ParagraphStyle(
        "Notes",
        parent=styles["Normal"],
        fontSize=9,
        leading=12,
        spaceAfter=8,
        textColor=HexColor("#333333"),
        leftIndent=10,
    )

    no_content_style = ParagraphStyle(
        "NoContent",
        parent=styles["Italic"],
        fontSize=9,
        textColor=HexColor("#999999"),
        spaceAfter=6,
        leftIndent=10,
    )

    # Build the document content
    story = []

    for idx, (image_path, note) in enumerate(
            zip(slide_images, notes), start=1
    ):
        # Add slide header
        story.append(Paragraph(f"Folie {idx}", slide_header_style))

        # Add slide image
        if image_path.exists():
            with Image.open(image_path) as img:
                img_width, img_height = img.size

            scale = page_width / img_width
            display_width = page_width
            display_height = img_height * scale

            slide_img = RLImage(
                str(image_path),
                width=display_width,
                height=display_height,
            )
            story.append(slide_img)
        else:
            story.append(Paragraph(
                f"<i>Folienbild nicht verfuegbar: {image_path.name}</i>",
                no_content_style
            ))

        story.append(Spacer(1, 0.3 * cm))

        # Add speaker notes
        story.append(Paragraph("Sprechernotizen:", notes_header_style))

        if note:
            formatted_note = escape_xml(note).replace("\n", "<br/>")
            story.append(Paragraph(formatted_note, notes_style))
        else:
            story.append(Paragraph(
                "Keine Sprechernotizen vorhanden.",
                no_content_style
            ))

        # Add page break after each slide (except the last one)
        if idx < len(slide_images):
            story.append(PageBreak())

    # Build the PDF
    doc.build(story)


def convert_pptx_to_pdf_with_notes(
        pptx_path: Path,
        output_path: Optional[Path] = None,
) -> Path:
    """
    Convert a PowerPoint presentation to PDF including speaker notes.

    Uses LibreOffice to render slides as images, then combines them
    with speaker notes extracted via python-pptx.

    Args:
        pptx_path: Path to the input PPTX file.
        output_path: Optional path for the output PDF.
                     If not provided, uses the same name as input with '.pdf' suffix.

    Returns:
        Path to the generated PDF file.

    Raises:
        FileNotFoundError: If input file or LibreOffice not found.
        RuntimeError: If conversion fails.
    """
    pptx_path = Path(pptx_path).resolve()

    if not pptx_path.exists():
        raise FileNotFoundError(f"Input file not found: {pptx_path}")

    # Find LibreOffice
    libreoffice_path = find_libreoffice()
    if not libreoffice_path:
        raise FileNotFoundError(
            "LibreOffice not found. Please install it:\n"
            "  macOS: brew install --cask libreoffice\n"
            "  Linux: sudo apt install libreoffice\n"
            "  Windows: Download from https://www.libreoffice.org/"
        )

    if output_path is None:
        output_path = pptx_path.parent / f"{pptx_path.stem}.pdf"
    else:
        output_path = Path(output_path).resolve()

    print(f"Converting: {pptx_path.name}")
    print(f"Output: {output_path}")
    print(f"Using LibreOffice: {libreoffice_path}")

    # Create temporary directory for intermediate files
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)

        # Step 1: Convert PPTX to PDF using LibreOffice
        print("Converting PPTX to PDF with LibreOffice...")
        cmd = [
            str(libreoffice_path),
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(temp_path),
            str(pptx_path),
        ]

        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=600,  # 10 minutes timeout
        )

        if result.returncode != 0:
            raise RuntimeError(
                f"LibreOffice conversion failed: {result.stderr}"
            )

        # Find the generated PDF
        pdf_file = temp_path / f"{pptx_path.stem}.pdf"
        if not pdf_file.exists():
            raise RuntimeError(
                f"PDF not generated. LibreOffice output: {result.stdout}"
            )

        print(f"PDF created: {pdf_file.name}")

        # Step 2: Extract PDF pages as images using PyMuPDF
        print("Extracting slide images from PDF...")
        pdf_doc = fitz.open(str(pdf_file))
        png_files = []

        # Use higher resolution for better quality (2x scale)
        zoom_matrix = fitz.Matrix(2.0, 2.0)

        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            pix = page.get_pixmap(matrix=zoom_matrix)

            # Save as PNG
            png_path = temp_path / f"slide_{page_num + 1:03d}.png"
            pix.save(str(png_path))
            png_files.append(png_path)

        pdf_doc.close()
        print(f"Extracted {len(png_files)} slide images")

        # Step 3: Extract speaker notes
        print("Extracting speaker notes...")
        notes = extract_speaker_notes(pptx_path)
        slide_count = get_slide_count(pptx_path)
        print(f"Found {slide_count} slides with {sum(1 for n in notes if n)} notes")

        # Verify we have the right number of images
        if len(png_files) != slide_count:
            print(
                f"Warning: Number of images ({len(png_files)}) "
                f"doesn't match slide count ({slide_count})"
            )
            # Pad notes if we have more images
            while len(notes) < len(png_files):
                notes.append("")
            # Truncate if we have fewer images
            notes = notes[:len(png_files)]

        # Step 4: Create final PDF with images and notes
        print("Creating PDF with slides and notes...")
        create_pdf_with_notes(png_files, notes, output_path)

    print(f"Successfully created: {output_path}")
    return output_path


def main() -> None:
    """Main entry point for the script."""
    # Define paths relative to project root
    script_dir = Path(__file__).parent
    project_root = script_dir.parent
    data_dir = project_root / "data"

    # Input PowerPoint file
    pptx_file = data_dir / "Inside Agentic AI.pptx"

    # Output PDF file
    output_file = data_dir / "Inside Agentic AI.pdf"

    try:
        convert_pptx_to_pdf_with_notes(pptx_file, output_file)
    except FileNotFoundError as e:
        print(f"Error: {e}")
        sys.exit(1)
    except Exception as e:
        print(f"Unexpected error: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
